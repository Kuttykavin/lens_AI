import cv2
import sys
import datetime
import tkinter as tk
from tkinter import messagebox, filedialog
from ultralytics import YOLO
import keyboard
import ctypes
import threading
import time
import os
import pystray
from PIL import Image, ImageDraw, ImageGrab, ImageTk
import queue
from pathlib import Path
import database

# ─── Config ───────────────────────────────────────────────────────────────────
CONF_THRESHOLD   = 0.3    # overridden by DB settings at runtime
MIN_BOX_RATIO    = 0.02
CONFIRM_FRAMES   = 2      # overridden by DB settings at runtime
LOCKDOWN_TIMEOUT = 30
DEBUG_PREVIEW    = True

THREAT_CLASSES = {67}   # 67 = cell phone (COCO)

# ─── Globals ──────────────────────────────────────────────────────────────────
protection_enabled = True
lockdown_active    = False
alert_queue        = queue.Queue()
_viewer            = None
_main_root         = None
_detection_model   = None
offense_count      = 0   # escalating lockdown counter

# Initialize database and evidence folder
database.init_database()
database.set_protection_status(True)  # mark active on startup
EVIDENCE_DIR = Path(".evidence")
EVIDENCE_DIR.mkdir(exist_ok=True)
if os.name == 'nt':  # Windows
    import ctypes
    ctypes.windll.kernel32.SetFileAttributesW(str(EVIDENCE_DIR), 2)  # Hidden

# ─── Tray icon ────────────────────────────────────────────────────────────────
def set_protection(icon, item):
    global protection_enabled
    protection_enabled = (str(item) == 'Enable Protection')
    database.set_protection_status(protection_enabled)

def toggle_preview(icon, item):
    global DEBUG_PREVIEW
    DEBUG_PREVIEW = not DEBUG_PREVIEW
    if not DEBUG_PREVIEW:
        cv2.destroyAllWindows()

def view_logs(icon, item):
    if os.path.exists('security_log.txt'):
        os.startfile('security_log.txt')

def create_icon_image():
    img  = Image.new('RGB', (64, 64), color='black')
    draw = ImageDraw.Draw(img)
    draw.polygon([(32,5),(10,15),(10,40),(32,60),(54,40),(54,15)], fill='blue')
    return img

menu = pystray.Menu(
    pystray.MenuItem('Enable Protection',   set_protection,  checked=lambda item: protection_enabled,     radio=True),
    pystray.MenuItem('Disable Protection',  set_protection,  checked=lambda item: not protection_enabled, radio=True),
    pystray.MenuItem('Show Camera Preview', toggle_preview,  checked=lambda item: DEBUG_PREVIEW),
    pystray.MenuItem('View Audit Logs',     view_logs)
)
tray_icon = pystray.Icon("ScreenSentry", create_icon_image(), "ScreenSentry", menu)
tray_icon.run_detached()

# ─── Camera tamper detection ──────────────────────────────────────────────────
def is_camera_tampered(frame):
    """Returns True if frame looks covered, blinded or smeared."""
    gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
    mean = gray.mean()

    # Fully black (covered with tape/hand)
    if mean < 8:
        return True, "Camera covered (black)"

    # Fully white/overexposed (flashlight)
    if mean > 247:
        return True, "Camera blinded (overexposed)"

    # Extremely blurry (smeared lens)
    blur = cv2.Laplacian(gray, cv2.CV_64F).var()
    if blur < 1.5:
        return True, "Camera obscured (blurry)"

    # Solid color object blocking camera (chips packet, paper, cloth, tape)
    # Low std deviation across all channels = uniform color = something blocking
    std_dev = gray.std()
    if std_dev < 18:
        return True, "Camera blocked (solid object)"

    # Check if a single color dominates >85% of the frame (e.g. orange chips packet)
    hsv = cv2.cvtColor(frame, cv2.COLOR_BGR2HSV)
    h, s, v = cv2.split(hsv)
    # High saturation + low variance in hue = solid colored object
    high_sat_pixels = (s > 60).sum()
    total_pixels    = frame.shape[0] * frame.shape[1]
    if high_sat_pixels / total_pixels > 0.75:
        hue_std = h[s > 60].std() if high_sat_pixels > 0 else 999
        if hue_std < 25:
            return True, "Camera blocked (colored object)"

    return False, ""

def handle_tamper(reason, frame):
    """Close document immediately and log tamper attempt."""
    print(f"[Tamper] {reason} — closing document")
    database.log_audit("Camera Tamper Detected", reason)
    # Save evidence
    try:
        timestamp  = datetime.datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
        image_path = str(EVIDENCE_DIR / f"tamper_{timestamp}.jpg")
        cv2.imwrite(image_path, frame)
        database.add_evidence(image_path, 1.0, 'High')
        database.send_threat_email(image_path, 1.0, 'High')
    except Exception as e:
        print(f"[Tamper] Evidence error: {e}")
    # Close document on main thread
    if _main_root:
        _main_root.after(0, close_document)
def capture_evidence(frame, confidence):
    """Save snapshot as evidence and log to database."""
    try:
        timestamp  = datetime.datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
        image_path = str(EVIDENCE_DIR / f"{timestamp}.jpg")
        cv2.imwrite(image_path, frame)

        # Determine severity based on confidence
        if confidence >= 0.75:
            severity = 'High'
        elif confidence >= 0.5:
            severity = 'Medium'
        else:
            severity = 'Low'

        database.add_evidence(image_path, round(confidence, 2), severity)
        print(f"[Evidence] Captured → {image_path} ({severity})")
    except Exception as e:
        print(f"[Evidence] Error capturing: {e}")


def anti_screenshot_loop():
    """Runs in background thread — uses the main root for clipboard ops."""
    while True:
        if protection_enabled and _main_root:
            try:
                img = ImageGrab.grabclipboard()
                if isinstance(img, Image.Image):
                    _main_root.after(0, lambda: (
                        _main_root.clipboard_clear(),
                        _main_root.update()
                    ))
            except Exception:
                pass
        time.sleep(0.5)

def screenshot_key_blocker():
    """Block PrintScreen and Win+Shift+S (Snipping Tool) hotkeys."""
    try:
        # Block PrintScreen
        keyboard.block_key('print screen')
        # Block Snipping Tool shortcut Win+Shift+S via suppress
        keyboard.add_hotkey('windows+shift+s', lambda: None, suppress=True)
        # Block Alt+PrintScreen
        keyboard.add_hotkey('alt+print screen', lambda: None, suppress=True)
        keyboard.wait()
    except Exception as e:
        print(f"[Blocker] Hotkey block error: {e}")

def screen_recorder_detector():
    """Detect if screen recording software is running."""
    RECORDERS = [
        'obs64.exe', 'obs32.exe', 'obs.exe',
        'bandicam.exe', 'fraps.exe', 'camtasia.exe',
        'screenrec.exe', 'sharex.exe', 'snagit32.exe',
        'flashback.exe', 'dxtory.exe', 'action.exe'
    ]
    while True:
        if protection_enabled:
            try:
                import psutil
                for proc in psutil.process_iter(['name']):
                    pname = proc.info['name'].lower()
                    if pname in RECORDERS:
                        print(f"[Security] Screen recorder detected: {pname}")
                        database.log_audit("Screen Recorder Detected", f"Process: {pname}")
                        alert_queue.put("LOCKDOWN")
                        break
            except ImportError:
                pass  # psutil not installed — skip
            except Exception:
                pass
        time.sleep(10)  # check every 10 seconds

threading.Thread(target=anti_screenshot_loop, daemon=True).start()

# ─── PPTX Viewer ──────────────────────────────────────────────────────────────
class PPTXViewer:
    def __init__(self, parent, path):
        self.slides = []
        self.index  = 0
        self.win    = tk.Toplevel(parent)
        self.win.title(f"ScreenSentry — {os.path.basename(path)}")
        self.win.configure(bg="black")
        self.win.protocol("WM_DELETE_WINDOW", self._on_close)

        self.canvas = tk.Canvas(self.win, bg="black", highlightthickness=0)
        self.canvas.pack(fill=tk.BOTH, expand=True)

        nav = tk.Frame(self.win, bg="#1a1a1a")
        nav.pack(fill=tk.X)

        self.btn_prev = tk.Button(nav, text="◀  Prev", command=self.prev_slide,
                                  bg="#333", fg="white", font=("Arial", 11),
                                  relief=tk.FLAT, padx=12, pady=4)
        self.btn_prev.pack(side=tk.LEFT, padx=8, pady=4)

        self.slide_label = tk.Label(nav, text="", bg="#1a1a1a", fg="white",
                                    font=("Arial", 11))
        self.slide_label.pack(side=tk.LEFT, expand=True)

        self.btn_next = tk.Button(nav, text="Next  ▶", command=self.next_slide,
                                  bg="#333", fg="white", font=("Arial", 11),
                                  relief=tk.FLAT, padx=12, pady=4)
        self.btn_next.pack(side=tk.RIGHT, padx=8, pady=4)

        self.win.bind("<Right>", lambda e: self.next_slide())
        self.win.bind("<Left>",  lambda e: self.prev_slide())

        self._load_pptx(path)

    def _load_pptx(self, path):
        loading_lbl = tk.Label(self.win, text="Loading slides...",
                               fg="white", bg="black", font=("Arial", 16))
        loading_lbl.place(relx=0.5, rely=0.5, anchor="center")
        self.win.update()

        def _convert():
            try:
                from pptx import Presentation
                from pptx.util import Pt
                import textwrap

                prs    = Presentation(path)
                width  = int(prs.slide_width  / 914400 * 96)
                height = int(prs.slide_height / 914400 * 96)
                slides_out = []
                for i, slide in enumerate(prs.slides):
                    W, H = max(width, 800), max(height, 600)
                    img  = Image.new("RGB", (W, H), "#1e1e2e")
                    draw = ImageDraw.Draw(img)
                    draw.rectangle([W-90, H-36, W-4, H-4], fill="#333355")
                    draw.text((W-86, H-32), f"{i+1} / {len(prs.slides)}", fill="#aaaacc")
                    y = 30
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if not line:
                                y += 10
                                continue
                            is_title = any(
                                run.font.size and run.font.size >= Pt(24)
                                for run in para.runs
                            )
                            color   = "#ffffff" if is_title else "#cccccc"
                            for wline in textwrap.wrap(line, width=70):
                                if y > H - 50:
                                    break
                                draw.text((40, y), wline, fill=color)
                                y += 28 if is_title else 22
                            y += 6
                    slides_out.append(img)
                self.slides = slides_out
            except ImportError:
                self.slides = [self._error_slide("python-pptx not installed.\nRun: pip install python-pptx")]
            except Exception as e:
                self.slides = [self._error_slide(f"Error loading:\n{e}")]

            loading_lbl.destroy()
            self.show_slide(0)

        threading.Thread(target=_convert, daemon=True).start()

    def _error_slide(self, msg):
        img  = Image.new("RGB", (800, 600), "#1e1e2e")
        draw = ImageDraw.Draw(img)
        draw.text((40, 260), msg, fill="#ff6666")
        return img

    def show_slide(self, idx):
        if not self.slides:
            return
        self.index    = max(0, min(idx, len(self.slides) - 1))
        slide_img     = self.slides[self.index].copy()
        cw = self.canvas.winfo_width()  or 900
        ch = self.canvas.winfo_height() or 600
        slide_img.thumbnail((cw, ch), Image.LANCZOS)
        self._tk_img = ImageTk.PhotoImage(slide_img)
        self.canvas.delete("all")
        self.canvas.create_image(cw//2, ch//2, anchor="center", image=self._tk_img)
        self.slide_label.config(text=f"Slide {self.index+1} of {len(self.slides)}")
        self.btn_prev.config(state=tk.NORMAL if self.index > 0 else tk.DISABLED)
        self.btn_next.config(state=tk.NORMAL if self.index < len(self.slides)-1 else tk.DISABLED)

    def next_slide(self): self.show_slide(self.index + 1)
    def prev_slide(self): self.show_slide(self.index - 1)
    def _on_close(self):  self.win.destroy()
    def hide(self):
        try: self.win.withdraw()
        except Exception: pass
    def show(self):
        try: self.win.deiconify(); self.win.lift()
        except Exception: pass
    def destroy(self):
        try: self.win.destroy()
        except Exception: pass

# ─── Document viewer ──────────────────────────────────────────────────────────
def open_document(path=None):
    global _viewer
    if path is None:
        path = filedialog.askopenfilename(
            title="Open Document to Protect",
            filetypes=[
                ("All supported", "*.pptx *.png *.jpg *.jpeg *.bmp *.gif *.txt"),
                ("PowerPoint",    "*.pptx"),
                ("Images",        "*.png *.jpg *.jpeg *.bmp *.gif"),
                ("Text files",    "*.txt"),
            ]
        )
    if not path:
        return

    if _viewer:
        try: _viewer.destroy()
        except Exception: pass
        _viewer = None

    ext = os.path.splitext(path)[1].lower()

    if ext == ".pptx":
        _viewer = PPTXViewer(_main_root, path)
        _viewer.win.geometry("960x620")
        return

    win = tk.Toplevel(_main_root)
    win.title(f"ScreenSentry — {os.path.basename(path)}")

    class _Generic:
        def __init__(self, w): self.win = w
        def hide(self):
            try: self.win.withdraw()
            except Exception: pass
        def show(self):
            try: self.win.deiconify(); self.win.lift()
            except Exception: pass
        def destroy(self):
            try: self.win.destroy()
            except Exception: pass

    _viewer = _Generic(win)

    if ext in (".png", ".jpg", ".jpeg", ".bmp", ".gif"):
        pil_img = Image.open(path)
        pil_img.thumbnail((900, 700))
        tk_img = ImageTk.PhotoImage(pil_img)
        lbl = tk.Label(win, image=tk_img)
        lbl.image = tk_img
        lbl.pack()
    else:
        frame = tk.Frame(win)
        frame.pack(fill=tk.BOTH, expand=True)
        sb  = tk.Scrollbar(frame)
        sb.pack(side=tk.RIGHT, fill=tk.Y)
        txt = tk.Text(frame, wrap=tk.WORD, yscrollcommand=sb.set, font=("Consolas", 11))
        txt.pack(fill=tk.BOTH, expand=True)
        sb.config(command=txt.yview)
        with open(path, 'r', errors='replace') as f:
            txt.insert(tk.END, f.read())
        txt.config(state=tk.DISABLED)

def close_document():
    global _viewer
    if _viewer:
        try: _viewer.destroy()
        except Exception: pass
        _viewer = None

def hide_document():
    if _viewer:
        try: _viewer.hide()
        except Exception: pass

def show_document():
    if _viewer:
        try: _viewer.show()
        except Exception: pass

# ─── Lockdown UI ──────────────────────────────────────────────────────────────
def trigger_defense():
    global lockdown_active, offense_count
    if lockdown_active:
        return
    lockdown_active = True
    offense_count  += 1
    hide_document()

    # ── Escalating: 3rd+ offense = immediate lock, no dialog ──────────────────
    if offense_count >= 3:
        print(f"[Lockdown] Offense #{offense_count} — immediate workstation lock")
        database.log_audit("Immediate Lockdown", f"Offense #{offense_count}")
        lockdown_active = False
        ctypes.windll.user32.LockWorkStation()
        return

    if not _main_root:
        print("[Lockdown] _main_root not ready — forcing workstation lock")
        ctypes.windll.user32.LockWorkStation()
        lockdown_active = False
        return

    # ── Escalating timeout: 1st=30s, 2nd=60s ──────────────────────────────────
    timeout = 30 if offense_count == 1 else 60

    try:
        overlay = tk.Toplevel(_main_root)
        overlay.attributes("-fullscreen", True)
        overlay.configure(bg="black")
        overlay.attributes("-topmost", True)
        overlay.overrideredirect(True)
        overlay.lift()
        overlay.focus_force()
    except Exception as e:
        print(f"[Lockdown] Overlay error: {e}")
        ctypes.windll.user32.LockWorkStation()
        lockdown_active = False
        return

    # ── Check if PIN is required ───────────────────────────────────────────────
    pin_enabled = database.is_lockdown_pin_enabled()
    correct_pin = database.get_lockdown_pin() if pin_enabled else ""

    countdown = [timeout]
    timer_id  = [None]
    resolved  = [False]

    # ── Countdown label ────────────────────────────────────────────────────────
    offense_label = tk.Label(
        overlay,
        text=f"⚠  Offense #{offense_count}  ⚠",
        fg="#ff6600", bg="black",
        font=("Arial", 16, "bold"), justify="center"
    )
    offense_label.pack(pady=(60, 0))

    label = tk.Label(
        overlay,
        text=f"⚠  Camera Detected  ⚠\n\nLocking in: {timeout}s",
        fg="red", bg="black", font=("Arial", 36, "bold"), justify="center"
    )
    label.pack(expand=True)

    def update_timer():
        global lockdown_active
        if resolved[0]:
            return
        countdown[0] -= 1
        if countdown[0] <= 0:
            resolved[0]     = True
            lockdown_active = False
            try: overlay.destroy()
            except Exception: pass
            ctypes.windll.user32.LockWorkStation()
        else:
            label.config(text=f"⚠  Camera Detected  ⚠\n\nLocking in: {countdown[0]}s")
            timer_id[0] = overlay.after(1000, update_timer)

    # ── PIN mode ───────────────────────────────────────────────────────────────
    if pin_enabled and correct_pin:
        pin_frame = tk.Frame(overlay, bg="black")
        pin_frame.pack(pady=10)

        pin_hint = tk.Label(
            pin_frame, text="Enter PIN to resume:",
            fg="#aaaacc", bg="black", font=("Arial", 14)
        )
        pin_hint.pack()

        pin_var = tk.StringVar()
        pin_entry = tk.Entry(
            pin_frame, textvariable=pin_var,
            show="●", font=("Arial", 24, "bold"),
            width=8, justify="center",
            bg="#1a1a2e", fg="white",
            insertbackground="white",
            relief=tk.FLAT, bd=0
        )
        pin_entry.pack(pady=8)
        pin_entry.focus_set()

        pin_msg = tk.Label(
            pin_frame, text="",
            fg="#ff3366", bg="black", font=("Arial", 11)
        )
        pin_msg.pack()

        def check_pin(event=None):
            global lockdown_active
            if resolved[0]:
                return
            entered = pin_var.get().strip()
            if entered == correct_pin:
                resolved[0]     = True
                lockdown_active = False
                print("[Lockdown] Correct PIN — resolved")
                if timer_id[0]:
                    try: overlay.after_cancel(timer_id[0])
                    except Exception: pass
                try: overlay.destroy()
                except Exception: pass
                show_document()
            else:
                pin_var.set("")
                pin_msg.config(text="❌ Incorrect PIN")
                pin_entry.config(bg="#2a0a14")
                overlay.after(800, lambda: pin_entry.config(bg="#1a1a2e"))

        pin_entry.bind("<Return>", check_pin)

        submit_btn = tk.Button(
            pin_frame, text="UNLOCK",
            command=check_pin,
            bg="#7b2fff", fg="white",
            font=("Arial", 12, "bold"),
            relief=tk.FLAT, padx=20, pady=6,
            cursor="hand2"
        )
        submit_btn.pack(pady=4)

    else:
        # ── No PIN — yes/no dialog ─────────────────────────────────────────────
        def ask_resume():
            global lockdown_active
            if resolved[0]:
                return
            try:
                answer = messagebox.askyesno(
                    "ScreenSentry — Security Alert",
                    f"Camera detected (Offense #{offense_count}).\n\nIs everything okay?\n\n"
                    "  Yes → Resume\n  No  → Close document",
                    parent=overlay
                )
            except Exception:
                answer = False
            if resolved[0]:
                return
            resolved[0]     = True
            lockdown_active = False
            print("[Lockdown] Resolved — lockdown_active reset to False")
            if timer_id[0]:
                try: overlay.after_cancel(timer_id[0])
                except Exception: pass
            try: overlay.destroy()
            except Exception: pass
            if answer:
                show_document()
            else:
                close_document()
                tray_icon.stop()
                os._exit(0)

        overlay.after(200, ask_resume)

    timer_id[0] = overlay.after(1000, update_timer)

# ─── YOLO detection thread ────────────────────────────────────────────────────
def detection_thread():
    model = YOLO("yolo26n.pt")

    print("\n=== Model Classes ===")
    for cid, cname in model.names.items():
        print(f"  {cid}: {cname}")
    print("=====================\n")

    cap = cv2.VideoCapture(0)
    cap.set(cv2.CAP_PROP_FRAME_WIDTH,  640)
    cap.set(cv2.CAP_PROP_FRAME_HEIGHT, 480)

    last_alert    = 0
    confirm_count = 0

    while True:
        if keyboard.is_pressed('ctrl+shift+q'):
            cap.release()
            cv2.destroyAllWindows()
            os._exit(0)

        success, img = cap.read()
        if not success:
            time.sleep(0.1)
            continue

        # ── Camera tamper check ──────────────────────────────────────────────
        if protection_enabled and _viewer:
            tampered, reason = is_camera_tampered(img)
            if tampered:
                threading.Thread(target=handle_tamper, args=(reason, img.copy()), daemon=True).start()
                time.sleep(3)  # avoid repeated triggers
                continue

        # Read sensitivity settings from DB dynamically
        det = database.get_detection_settings()
        conf_threshold  = det["conf_threshold"]
        confirm_frames  = det["confirm_frames"]

        frame_h, frame_w = img.shape[:2]
        phone_this_frame = False
        display_img      = img.copy()

        if protection_enabled:
            results = model.predict(img, stream=True, verbose=False, conf=conf_threshold)
            for r in results:
                for box in r.boxes:
                    class_id   = int(box.cls[0])
                    conf_score = float(box.conf[0])
                    x1, y1, x2, y2 = map(int, box.xyxy[0])
                    bw = x2 - x1
                    bh = y2 - y1

                    too_small = (bw / frame_w < MIN_BOX_RATIO or
                                 bh / frame_h < MIN_BOX_RATIO)
                    is_phone  = (class_id in THREAT_CLASSES and not too_small)

                    if DEBUG_PREVIEW:
                        color     = (0, 0, 255) if is_phone else (0, 200, 0)
                        label_txt = (f"PHONE {conf_score:.2f}" if is_phone
                                     else f"cls:{class_id} {conf_score:.2f}")
                        cv2.rectangle(display_img, (x1, y1), (x2, y2), color, 2)
                        cv2.putText(display_img, label_txt, (x1, max(y1-8, 10)),
                                    cv2.FONT_HERSHEY_SIMPLEX, 0.55, color, 2)
                        print(f"  Detected → cls:{class_id} ({model.names.get(class_id,'?')}) conf:{conf_score:.2f} box:({bw}x{bh}) too_small:{too_small}")

                    if is_phone:
                        phone_this_frame = True

        if phone_this_frame:
            confirm_count += 1
        else:
            confirm_count = 0

        threat = (confirm_count >= confirm_frames)

        if DEBUG_PREVIEW:
            if confirm_count > 0 and not threat:
                status_txt, status_color = f"Confirming... {confirm_count}/{confirm_frames}", (0, 165, 255)
            elif threat:
                status_txt, status_color = "THREAT CONFIRMED", (0, 0, 255)
            else:
                status_txt, status_color = "Monitoring...", (0, 255, 0)
            cv2.putText(display_img, status_txt, (10, 30),
                        cv2.FONT_HERSHEY_SIMPLEX, 0.9, status_color, 2)
            cv2.imshow("ScreenSentry — Camera Feed", display_img)
            cv2.waitKey(1)

        now = time.time()
        if threat and (now - last_alert) > 2:
            last_alert    = now
            confirm_count = 0
            print("[Lockdown] THREAT CONFIRMED — queuing lockdown")

            # Capture evidence snapshot
            best_conf = 0.0
            img_copy  = img.copy()
            if protection_enabled:
                for r in model.predict(img, stream=True, verbose=False, conf=conf_threshold):
                    for box in r.boxes:
                        if int(box.cls[0]) in THREAT_CLASSES:
                            best_conf = max(best_conf, float(box.conf[0]))

            # Determine severity
            if best_conf >= 0.75:
                severity = 'High'
            elif best_conf >= 0.5:
                severity = 'Medium'
            else:
                severity = 'Low'

            # Save evidence and send email alert in background
            def _handle_threat(frame, conf, sev):
                capture_evidence(frame, conf)
                time.sleep(0.5)
                try:
                    files = sorted(EVIDENCE_DIR.glob("*.jpg"), key=os.path.getmtime, reverse=True)
                    img_path = str(files[0]) if files else None
                    print(f"[Alert] Sending threat email — {sev} {conf:.2f} img:{img_path}")
                    database.send_threat_email(img_path, conf, sev)
                except Exception as e:
                    print(f"[Alert] Email error: {e}")

            threading.Thread(target=_handle_threat, args=(img_copy, best_conf, severity), daemon=True).start()

            with open('security_log.txt', 'a') as f:
                f.write(
                    f"[{datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] "
                    f"[ALERT] Camera Detected — YOLOv26 Verified\n"
                )
            # Only trigger lockdown UI if not already in lockdown
            if not lockdown_active:
                print(f"[Lockdown] Putting LOCKDOWN in queue, lockdown_active={lockdown_active}")
                alert_queue.put("LOCKDOWN")
                print(f"[Lockdown] Queue size now: {alert_queue.qsize()}")

        time.sleep(0.005)

# ─── Main UI ──────────────────────────────────────────────────────────────────
def ui_thread():
    global _main_root
    _main_root = tk.Tk()
    _main_root.title("ScreenSentry")
    _main_root.geometry("420x220")
    _main_root.resizable(False, False)
    _main_root.configure(bg="#0f0f1a")

    tk.Label(_main_root, text="🛡  ScreenSentry", font=("Arial", 20, "bold"),
             bg="#0f0f1a", fg="white").pack(pady=10)
    tk.Label(_main_root, text="Anti-Screenshot Protection Active",
             font=("Arial", 11), bg="#0f0f1a", fg="#00cc66").pack()

    tk.Button(
        _main_root, text="📂  Open Document to Protect",
        command=lambda: open_document(),
        font=("Arial", 11), bg="#2196F3", fg="white",
        relief=tk.FLAT, padx=10, pady=6, cursor="hand2"
    ).pack(pady=20)

    tk.Label(_main_root, text="Supports: .pptx  .png  .jpg  .txt",
             font=("Arial", 9), bg="#0f0f1a", fg="#666688").pack()
    tk.Label(_main_root, text="Press Ctrl+Shift+Q to quit",
             font=("Arial", 9), bg="#0f0f1a", fg="#444466").pack(pady=4)

    def poll_alerts():
        if not alert_queue.empty():
            msg = alert_queue.get()
            print(f"[Lockdown] poll_alerts got message: {msg}")
            if msg == "LOCKDOWN":
                print("[Lockdown] Calling trigger_defense NOW")
                _main_root.after(0, trigger_defense)
        _main_root.after(50, poll_alerts)

    _main_root.after(100, poll_alerts)

    try:
        _main_root.mainloop()
    finally:
        tray_icon.stop()

# ─── Entry point ──────────────────────────────────────────────────────────────
if __name__ == "__main__":
    threading.Thread(target=detection_thread,       daemon=True).start()
    threading.Thread(target=screenshot_key_blocker, daemon=True).start()
    threading.Thread(target=screen_recorder_detector, daemon=True).start()

    if len(sys.argv) > 1 and os.path.exists(sys.argv[1]):
        def _delayed_open():
            time.sleep(0.8)
            if _main_root:
                _main_root.after(0, lambda: open_document(sys.argv[1]))
        threading.Thread(target=_delayed_open, daemon=True).start()

    ui_thread()
